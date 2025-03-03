import re

import pdfplumber
from django.core.files.storage import default_storage
from pptx import Presentation
from rest_framework.generics import ListAPIView, RetrieveAPIView
from rest_framework.parsers import MultiPartParser, FormParser
from rest_framework.response import Response
from rest_framework.views import APIView

from apps.models import Question, TestQuestion


def extract_questions_from_pdf(pdf_path):
    """ PDF fayldan test savollarini ajratib olish """
    questions = []
    with pdfplumber.open(pdf_path) as pdf:
        text = "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())

    pattern = re.compile(r"(\d+\)\s.*?\?)\s*(a\..*?)\s*(b\..*?)\s*(c\..*?)\s*(d\..*?)", re.DOTALL)
    matches = pattern.findall(text)

    for match in matches:
        question = match[0]
        options = match[1:]
        questions.append({
            "question": question.strip(),
            "options": [opt.strip() for opt in options],
            "correct_answer": options[0]  # Birinchi variantni default to'g'ri deb olish
        })

    return questions


def extract_questions_from_pptx(pptx_path):
    """ PPTX fayldan test savollarini ajratib olish """
    prs = Presentation(pptx_path)
    questions = []

    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]

        if len(texts) >= 2:
            question = texts[0]  # Birinchi qator - savol
            answers = texts[1:]  # Variantlar

            if len(answers) >= 2:
                correct_answer = answers[0]  # Default birinchi javob to'g'ri deb olinadi
                questions.append({
                    "question": question.strip(),
                    "options": [opt.strip() for opt in answers],
                    "correct_answer": correct_answer
                })

    return questions


def save_test_to_db(file_path, test_file):
    """ PDF yoki PPTX fayldan testlarni chiqarib bazaga yozish """
    file_extension = file_path.split(".")[-1].lower()

    if file_extension == "pdf":
        test_data = extract_questions_from_pdf(file_path)
    elif file_extension == "pptx":
        test_data = extract_questions_from_pptx(file_path)
    else:
        raise ValueError("❌ Faqat PDF yoki PPTX fayllar qo‘llab-quvvatlanadi!")

    for test in test_data:
        Question.objects.create(
            file=test_file,
            text=test["question"],
            options=test["options"],
            correct_answer=test["options"][0]
        )

    print(f"✅ {len(test_data)} ta test bazaga saqlandi!")


class UploadTestFileView(APIView):
    parser_classes = (MultiPartParser, FormParser)

    def post(self, request, *args, **kwargs):
        file = request.FILES.get('file')
        test_file = request.date.get('test_file')
        if not file:
            return Response({"error": "Fayl yuklanmadi!"}, status=400)

        file_path = default_storage.save(f"uploads/{file.name}", file, test_file)
        full_path = default_storage.path(file_path, test_file)

        try:
            save_test_to_db(full_path)
            return Response({"message": "✅ Testlar bazaga saqlandi!"}, status=201)
        except Exception as e:
            return Response({"error": str(e)}, status=400)


class TestFileListView(ListAPIView):
    """Barcha yuklangan fayllarni testlari bilan chiqarish"""
    queryset = TestQuestion.objects.all()
    serializer_class = TestFileSerializer


class TestFileDetailView(RetrieveAPIView):
    """Bitta faylga tegishli testlarni chiqarish"""
    queryset = TestQuestion.objects.all()
    serializer_class = TestFileSerializer
    lookup_field = "id"
